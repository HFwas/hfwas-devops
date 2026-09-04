#!/usr/bin/env python3
"""
文档生成脚本 - 最简单的文档生成方式
支持：Word / Excel / PPT / 图片 / Markdown / PDF

用法：
  python3 generate_doc.py <format> <data.json> <output_path>

示例：
  python3 generate_doc.py word data.json output.docx
  python3 generate_doc.py excel data.json output.xlsx
  python3 generate_doc.py ppt  data.json output.pptx
  python3 generate_doc.py image data.json output.png
  python3 generate_doc.py md   data.json output.md
  python3 generate_doc.py pdf  data.json output.pdf

data.json 格式见各生成函数注释。
"""

import json
import sys
import os
import io
import random
import string
import platform
import subprocess
from datetime import datetime
import math


def _as_bool(value) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value != 0
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "on"}
    return False


# ============================================================
# 自动配置中文字体（matplotlib 图表中文显示）
# ============================================================
def _setup_chinese_font():
    """自动检测并配置 matplotlib 中文字体"""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.font_manager import FontProperties, fontManager

        # 常见中文字体路径（按优先级）
        font_candidates = [
            # macOS
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
            # Windows
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simhei.ttf",
            "C:/Windows/Fonts/simsun.ttc",
            # Linux
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        ]

        for font_path in font_candidates:
            if os.path.exists(font_path):
                fontManager.addfont(font_path)
                font_name = FontProperties(fname=font_path).get_name()
                plt.rcParams["font.family"] = font_name
                plt.rcParams["axes.unicode_minus"] = False
                return

        # 如果没找到，尝试从系统查找
        result = subprocess.run(
            ["fc-list", ":lang=zh", "-f", "%{file}\n"],
            capture_output=True, text=True, timeout=5
        )
        if result.returncode == 0:
            for line in result.stdout.strip().split("\n"):
                line = line.strip()
                if line and os.path.exists(line):
                    fontManager.addfont(line)
                    font_name = FontProperties(fname=line).get_name()
                    plt.rcParams["font.family"] = font_name
                    plt.rcParams["axes.unicode_minus"] = False
                    return

    except Exception:
        pass  # 中文不可用就 fallback 到英文


_setup_chinese_font()


# ============================================================
# 模拟数据生成器（Excel 专用）
# ============================================================
import random
import string

_MOCK_FIRST_NAMES = [
    "王", "李", "张", "刘", "陈", "杨", "黄", "赵", "周", "吴",
    "徐", "孙", "马", "胡", "朱", "郭", "何", "罗", "高", "林",
]
_MOCK_LAST_NAMES = [
    "伟", "芳", "娜", "秀英", "敏", "静", "丽", "强", "磊", "军",
    "洋", "勇", "艳", "杰", "娟", "涛", "明", "超", "秀兰", "霞",
    "平", "刚", "桂英", "文", "华", "飞", "玉兰", "斌", "玲", "国",
]
_MOCK_DEPARTMENTS = [
    "技术部", "市场部", "销售部", "财务部", "人事部", "研发部",
    "产品部", "运营部", "客服部", "行政部", "测试部", "设计部",
]
_MOCK_POSITIONS = [
    "工程师", "经理", "主管", "总监", "专员", "实习生",
    "高级工程师", "架构师", "分析师", "顾问",
]
_MOCK_STATUSES = ["进行中", "已完成", "待开始", "已暂停", "已关闭"]
_MOCK_CITIES = [
    "北京", "上海", "广州", "深圳", "杭州", "成都", "武汉", "南京",
    "西安", "重庆", "苏州", "天津", "长沙", "郑州", "东莞", "青岛",
]
_MOCK_PROJECTS = [
    "Alpha项目", "Beta平台", "Gamma系统", "Delta升级",
    "数据中台", "用户中心", "订单系统", "支付平台",
]


def _mock_name():
    return random.choice(_MOCK_FIRST_NAMES) + random.choice(_MOCK_LAST_NAMES)


def _mock_phone():
    prefixes = ["138", "139", "150", "151", "152", "186", "187", "188"]
    return random.choice(prefixes) + "".join(random.choices(string.digits, k=8))


def _mock_email(name=None):
    domains = ["qq.com", "163.com", "gmail.com", "outlook.com", "company.cn"]
    local = (name or _mock_name()).lower() + str(random.randint(1, 999))
    return f"{local}@{random.choice(domains)}"


def _mock_date(start_year=2024, end_year=2026):
    y = random.randint(start_year, end_year)
    m = random.randint(1, 12)
    d = random.randint(1, 28)
    return f"{y}-{m:02d}-{d:02d}"


def _mock_datetime():
    return f"{_mock_date()} {random.randint(9, 18):02d}:{random.randint(0, 59):02d}"


def _mock_address():
    city = random.choice(_MOCK_CITIES)
    district = random.choice(["朝阳区", "海淀区", "天河区", "武侯区", "西湖区",
                               "浦东新区", "洪山区", "鼓楼区"])
    street = f"{random.choice(['中山', '人民', '解放', '建设', '科技'])}路{random.randint(1, 999)}号"
    return f"{city}{district}{street}"


def _mock_id_card():
    area = random.choice(["110101", "310101", "440101", "320101", "510101"])
    birth = f"{random.randint(1970, 2005)}{random.randint(1, 12):02d}{random.randint(1, 28):02d}"
    seq = "".join(random.choices(string.digits, k=3))
    return area + birth + seq + random.choice(string.digits + "X")


def _mock_salary():
    return random.randint(5000, 50000)


def _mock_age():
    return random.randint(22, 55)


def _mock_score():
    return round(random.uniform(60, 100), 1)


def _mock_progress():
    return f"{random.randint(0, 100)}%"


def _guess_column_type(col_name: str) -> str:
    """根据列名猜测数据类型"""
    name = col_name.lower()
    if any(k in name for k in ["姓名", "名字", "名称", "name", "员工", "人员", "用户"]):
        return "name"
    if any(k in name for k in ["手机", "电话", "手机号", "联系电话", "phone", "tel"]):
        return "phone"
    if any(k in name for k in ["邮箱", "邮件", "email", "mail"]):
        return "email"
    if any(k in name for k in ["年龄", "age"]):
        return "age"
    if any(k in name for k in ["日期", "date", "时间", "time", "创建", "更新"]):
        return "datetime" if any(k in name for k in ["时间", "time"]) else "date"
    if any(k in name for k in ["地址", "address", "住址", "位置"]):
        return "address"
    if any(k in name for k in ["身份证", "id", "证件"]):
        return "id_card"
    if any(k in name for k in ["薪资", "工资", "薪水", "salary", "收入", "薪酬"]):
        return "salary"
    if any(k in name for k in ["部门", "dept", "depart", "团队"]):
        return "department"
    if any(k in name for k in ["职位", "岗位", "职务", "position", "title", "职称"]):
        return "position"
    if any(k in name for k in ["状态", "status", "阶段"]):
        return "status"
    if any(k in name for k in ["城市", "city", "地区"]):
        return "city"
    if any(k in name for k in ["分数", "score", "成绩", "评分"]):
        return "score"
    if any(k in name for k in ["进度", "progress", "完成"]):
        return "progress"
    if any(k in name for k in ["项目", "project"]):
        return "project"
    if any(k in name for k in ["金额", "price", "价格", "费用", "成本", "amount", "总额"]):
        return "salary"
    return "text"


_MOCK_GENERATORS = {
    "name": lambda: _mock_name(),
    "phone": lambda: _mock_phone(),
    "email": lambda: _mock_email(),
    "age": lambda: _mock_age(),
    "date": lambda: _mock_date(),
    "datetime": lambda: _mock_datetime(),
    "address": lambda: _mock_address(),
    "id_card": lambda: _mock_id_card(),
    "salary": lambda: _mock_salary(),
    "department": lambda: random.choice(_MOCK_DEPARTMENTS),
    "position": lambda: random.choice(_MOCK_POSITIONS),
    "status": lambda: random.choice(_MOCK_STATUSES),
    "city": lambda: random.choice(_MOCK_CITIES),
    "score": lambda: _mock_score(),
    "progress": lambda: _mock_progress(),
    "project": lambda: random.choice(_MOCK_PROJECTS),
    "text": lambda: f"数据{random.randint(1, 9999)}",
}


def generate_mock_rows(columns: list, row_count: int = 10, row_size: int = 10) -> list:
    """
    根据列名自动生成模拟数据行

    columns: 列名列表，如 ["姓名", "年龄", "部门", "薪资"]
             或 dict 列表，如 [{"name": "姓名", "type": "name"}, ...]
    row_count: 生成行数
    row_size: 每个单元格的字符数/数据量（仅对字符串类型生效）

    返回: [["王伟", 28, "技术部", 15000], ...]
    """
    col_configs = []
    for col in columns:
        if isinstance(col, dict):
            col_configs.append((col.get("name", ""), col.get("type", "text")))
        else:
            name = str(col)
            col_type = _guess_column_type(name)
            col_configs.append((name, col_type))

    rows = []
    for _ in range(row_count):
        row = []
        for col_name, col_type in col_configs:
            if col_type == "text" and row_size > 10:
                # 根据 row_size 生成对应长度的随机文本
                # 每 10 个 row_size 生成约 1 个中文字符的等效内容
                text_len = max(1, row_size // 3)
                chars = string.ascii_letters + string.digits + " "
                row.append(''.join(random.choices(chars, k=text_len)))
            elif col_type == "name" and row_size > 10:
                # 姓名可以加长（双名变多名）
                extra = max(0, (row_size - 10) // 5)
                row.append(_mock_name() + ''.join(random.choices(['伟','芳','娜','强','磊','军','洋','勇','杰','涛'], k=extra)))
            elif col_type == "address" and row_size > 10:
                # 地址加长
                extra = f"第{random.randint(1, row_size)}号院{random.randint(1, row_size)}号楼"
                row.append(_mock_address() + extra)
            elif col_type == "email" and row_size > 10:
                # 邮箱加长
                local_len = max(5, row_size // 3)
                local = ''.join(random.choices(string.ascii_lowercase + string.digits, k=local_len))
                row.append(f"{local}@{random.choice(['company.cn', 'corp.com', 'group.net'])}")
            else:
                generator = _MOCK_GENERATORS.get(col_type, _MOCK_GENERATORS["text"])
                row.append(generator())
        rows.append(row)

    return rows

# ============================================================
# Word
# ============================================================
def generate_word(data: dict, output_path: str):
    """
    data 格式:
    {
        "title": "标题",
        "paragraphs": ["段落1", "段落2"],
        "table": {
            "headers": ["列1", "列2"],
            "rows": [["a1", "b1"], ["a2", "b2"]]
        }
    }
    """
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # 标题
    doc.add_heading(data.get("title", "文档"), level=0)

    # 段落
    for p in data.get("paragraphs", []):
        if isinstance(p, str):
            doc.add_paragraph(p)
        elif isinstance(p, dict):
            para = doc.add_paragraph(p.get("text", ""))
            if p.get("bold"):
                for run in para.runs:
                    run.bold = True
            if p.get("align") == "center":
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 表格
    table_data = data.get("table")
    if table_data:
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Light Grid Accent 1"
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
        for r_idx, row in enumerate(rows):
            for c_idx, cell in enumerate(row):
                table.rows[r_idx + 1].cells[c_idx].text = str(cell)

    # 图片
    img_path = data.get("image")
    if img_path and os.path.exists(img_path):
        doc.add_picture(img_path, width=Inches(5.5))

    doc.save(output_path)
    print(f"[OK] Word 文档已生成: {output_path}")


# ============================================================
# Excel
# ============================================================
def generate_excel(data: dict, output_path: str):
    """
    data 格式:
    {
        "sheet_name": "Sheet1",

        # 方式一：手动数据
        "rows": [
            ["姓名", "年龄", "部门"],
            ["张三", 28, "技术部"],
            ["李四", 32, "市场部"]
        ],

        # 方式二：自动生成模拟数据
        "mock": true,
        "columns": ["姓名", "年龄", "部门", "薪资", "手机号"],
        "row_count": 20
    }
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = data.get("sheet_name", "Sheet1")

    # 生成数据
    rows = data.get("rows", [])
    if data.get("mock"):
        columns = data.get("columns", [])
        row_count = data.get("row_count", 10)
        # 自动生成
        generated = generate_mock_rows(columns, row_count)
        # 第一行是表头
        headers = [c["name"] if isinstance(c, dict) else str(c) for c in columns]
        rows = [headers] + generated

    # 写入数据
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row, start=1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    # 表头样式
    if data.get("header_style", True) and rows:
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4",
                                  end_color="4472C4",
                                  fill_type="solid")
        for cell in ws[1]:
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")

    # 自动列宽
    for col in ws.columns:
        max_len = max(len(str(c.value or "")) for c in col)
        ws.column_dimensions[get_column_letter(col[0].column)].width = max_len + 2

    wb.save(output_path)
    print(f"[OK] Excel 文件已生成: {output_path}")


def generate_excel_mock(data: dict, output_path: str):
    """自动生成模拟数据的 Excel 文件"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    from openpyxl.utils import get_column_letter

    columns = data.get("columns", [])
    column_count = data.get("column_count", 0)
    if column_count > 0:
        # column_count 优先：自动生成列名
        columns = [f"列{i}" for i in range(1, column_count + 1)]
    elif not columns:
        # 都没有设置时，沿用默认列
        columns = ["姓名", "年龄", "部门", "薪资", "手机号", "邮箱", "入职日期", "状态"]
    row_count = data.get("row_count", 10)
    row_size = data.get("row_size", 10)
    file_count = data.get("file_count", 1)
    output_dir = data.get("output_dir", "")
    headers = [c["name"] if isinstance(c, dict) else str(c) for c in columns]

    def _write_one(path: str, data_rows: list) -> int:
        wb = Workbook()
        ws = wb.active
        ws.title = data.get("sheet_name", "Sheet1")
        for ri, row in enumerate(data_rows, start=1):
            for ci, val in enumerate(row, start=1):
                ws.cell(row=ri, column=ci, value=val)
        if data_rows:
            hf = Font(bold=True, color="FFFFFF")
            hfill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            for cell in ws[1]:
                cell.font = hf
                cell.fill = hfill
                cell.alignment = Alignment(horizontal="center")
        for col in ws.columns:
            ml = max(len(str(c.value or "")) for c in col)
            ws.column_dimensions[get_column_letter(col[0].column)].width = ml + 2
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        wb.save(path)
        return os.path.getsize(path)

    if file_count <= 1:
        generated = generate_mock_rows(columns, row_count, row_size=row_size)
        save_path = output_path if not output_dir else os.path.join(output_dir, os.path.basename(output_path))
        _write_one(save_path, [headers] + generated)
        size = os.path.getsize(save_path)
        print(f"[OK] Excel 文件已生成: {save_path} ({size} bytes)")
    else:
        rows_per_file = max(1, row_count // file_count)
        base = os.path.splitext(output_path)[0]
        base_name = os.path.basename(base)
        ext = os.path.splitext(output_path)[1]
        total_size = 0
        files = []
        for i in range(file_count):
            start = i * rows_per_file
            end = start + rows_per_file if i < file_count - 1 else row_count
            generated = generate_mock_rows(columns, end - start, row_size=row_size)
            data_rows = [headers] + generated
            part_name = f"{base_name}_{i + 1}{ext}"
            if output_dir:
                part_path = os.path.join(output_dir, part_name)
            else:
                part_path = f"{base}_{i + 1}{ext}"
            # 第一个文件也写入原始 output_path，浏览器下载时能读到
            if i == 0 and not output_dir:
                _write_one(output_path, data_rows)
            _write_one(part_path, data_rows)
            size = os.path.getsize(part_path)
            total_size += size
            files.append(part_path)
            print(f"  [{i + 1}/{file_count}] {os.path.basename(part_path)} ({size} bytes)")
        print(f"[OK] Excel {file_count} 个文件已生成，总大小: {total_size} bytes")


# ============================================================
# PPT
# ============================================================
def generate_ppt(data: dict, output_path: str):
    """
    data 格式:
    {
        "title": "演示文稿标题",
        "subtitle": "副标题",
        "slides": [
            {"title": "第一页", "items": ["内容1", "内容2"]},
            {"title": "第二页", "items": ["内容A", "内容B"]}
        ]
    }
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get("title", "演示文稿")
    subtitle = data.get("subtitle", "")
    if subtitle:
        slide.placeholders[1].text = subtitle

    # 内容页
    for page in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = page.get("title", "")
        content = slide.placeholders[1]
        content.text = "\n".join(page.get("items", []))

    prs.save(output_path)
    print(f"[OK] PPT 文件已生成: {output_path}")


# ============================================================
# 图片
# ============================================================
def generate_image(data: dict, output_path: str):
    """
    data 格式:
    {
        "chart_type": "line|bar|pie",
        "title": "图表标题",
        # 折线图
        "series": [
            {"name": "系列1", "x": ["1月","2月"], "y": [100, 150]},
            {"name": "系列2", "x": ["1月","2月"], "y": [200, 180]}
        ],
        # 柱状图
        "categories": ["A", "B", "C"],
        "values": [10, 20, 15],
        # 饼图
        "labels": ["分类A", "分类B"],
        "sizes": [60, 40],
        # 纯图片
        "width": 800, "height": 400, "color": "white"
    }
    """
    import matplotlib.pyplot as plt
    from PIL import Image

    chart_type = data.get("chart_type", "line")
    # 支持从 data 中指定 figsize 和 dpi（用于文件大小调整）
    figsize = data.get("figsize", (10, 6))
    dpi = data.get("dpi", 150)

    if chart_type == "line":
        fig, ax = plt.subplots(figsize=figsize)
        for series in data.get("series", []):
            ax.plot(series.get("x", []), series.get("y", []),
                    label=series.get("name", ""), marker="o")
        ax.set_title(data.get("title", "折线图"))
        ax.legend()
        ax.grid(True, alpha=0.3)
        fig.tight_layout()
        fig.savefig(output_path, dpi=150)
        plt.close()

    elif chart_type == "bar":
        fig, ax = plt.subplots(figsize=figsize)
        ax.bar(data.get("categories", []), data.get("values", []),
               color="steelblue")
        ax.set_title(data.get("title", "柱状图"))
        fig.tight_layout()
        fig.savefig(output_path, dpi=dpi)
        plt.close()

    elif chart_type == "pie":
        fig, ax = plt.subplots(figsize=figsize)
        ax.pie(data.get("sizes", []), labels=data.get("labels", []),
               autopct="%1.1f%%", startangle=90)
        ax.axis("equal")
        ax.set_title(data.get("title", "饼图"))
        fig.savefig(output_path, dpi=dpi)
        plt.close()

    elif chart_type == "scatter":
        import numpy as np
        np.random.seed(42)
        fig, ax = plt.subplots(figsize=figsize)
        n_points = data.get("n_points", 1000)
        x = np.random.rand(n_points)
        y = np.random.rand(n_points)
        ax.scatter(x, y, s=1, alpha=0.5)
        ax.set_title(data.get("title", "散点图"))
        fig.savefig(output_path, dpi=dpi, bbox_inches="tight")
        plt.close()

    else:
        # 纯图片
        img = Image.new("RGB",
                        (data.get("width", 800), data.get("height", 400)),
                        color=data.get("color", "white"))
        img.save(output_path)

    print(f"[OK] 图片已生成: {output_path}")


# ============================================================
# Markdown
# ============================================================
def generate_markdown(data: dict, output_path: str):
    """
    data 格式:
    {
        "title": "标题",
        "date": "2026-08-27",
        "author": "作者",
        "sections": [
            {"heading": "章节1", "content": "章节内容"},
            {"heading": "章节2", "content": "章节内容"}
        ],
        "table": {
            "headers": ["列1", "列2"],
            "rows": [["a1", "b1"], ["a2", "b2"]]
        }
    }
    """
    md = []
    md.append(f"# {data.get('title', '文档')}")
    md.append("")

    date = data.get("date", "")
    author = data.get("author", "")
    if date or author:
        md.append("> " + ", ".join(filter(None, [date, author])))
        md.append("")

    # 章节
    for section in data.get("sections", []):
        md.append(f"## {section.get('heading', '')}")
        md.append("")
        md.append(section.get("content", ""))
        md.append("")

    # 表格
    table = data.get("table")
    if table:
        headers = table.get("headers", [])
        rows = table.get("rows", [])
        md.append("| " + " | ".join(headers) + " |")
        md.append("| " + " | ".join("---" for _ in headers) + " |")
        for row in rows:
            md.append("| " + " | ".join(str(c) for c in row) + " |")
        md.append("")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md))

    print(f"[OK] Markdown 文件已生成: {output_path}")


# ============================================================
# PDF
# ============================================================
def generate_pdf(data: dict, output_path: str):
    """
    data 格式:
    {
        "title": "文档标题",
        "content": ["段落1", "段落2", "段落3"],
        "table": {
            "headers": ["列1", "列2"],
            "rows": [["a1", "b1"], ["a2", "b2"]]
        },
        "page_count": 5,          # 可选，指定生成页数
        "empty_page_count": false, # true 时不指定页数，按内容自动分页
        "empty_content": false,    # true 时不写正文（空白页）
        "encrypt": false,          # true 时加密
        "pdf_password": "123456"
    }
    """
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.pdfencrypt import StandardEncryption

    empty_content = _as_bool(data.get("empty_content"))
    empty_page_count = _as_bool(data.get("empty_page_count"))
    encrypt = _as_bool(data.get("encrypt"))
    password = str(data.get("pdf_password") or data.get("password") or "123456")

    page_count = 0 if empty_page_count else int(data.get("page_count") or 0)
    content = [] if empty_content else list(data.get("content") or [])
    if not empty_content and not content:
        content = ["这是第一段内容，由 Python 脚本自动生成。", "这是第二段内容。"]
    total_paragraphs = len(content)
    title = data.get("title", "文档")

    encrypt_obj = None
    if encrypt:
        encrypt_obj = StandardEncryption(userPassword=password, ownerPassword=password)

    width, height = A4
    c = canvas.Canvas(output_path, pagesize=A4, encrypt=encrypt_obj)

    if empty_content:
        blank_pages = page_count if page_count > 0 else 1
        for i in range(blank_pages):
            if i > 0:
                c.showPage()
        c.save()
        flags = []
        if encrypt:
            flags.append("加密")
        flags.append("空内容")
        if empty_page_count:
            flags.append("空页数")
        print(f"[OK] PDF 文件已生成: {output_path}，共 {blank_pages} 页（{', '.join(flags)}）")
        return

    # A4: 595×842 pt, 可用高度约 660pt
    # 每行 8mm ≈ 22.7pt, 每页约 29 行
    # 每段落约 5 行（含自动换行），每页约 5-6 段落
    MAX_PARAS_PER_PAGE = 6
    LINE_HEIGHT = 8 * mm

    if page_count > 0 and total_paragraphs > 0:
        # 按指定页数均分段落
        paras_per_page = math.ceil(total_paragraphs / page_count)
    else:
        paras_per_page = MAX_PARAS_PER_PAGE

    para_index = 0
    page_num = 0

    while para_index < total_paragraphs or (page_count > 0 and page_num < page_count):
        page_num += 1
        need_content = page_count > 0 and page_num > 1 and para_index >= total_paragraphs

        if page_num == 1:
            c.setFont("Helvetica-Bold", 24)
            c.drawString(50 * mm, height - 30 * mm, title)
            y = height - 50 * mm
        elif need_content:
            c.setFont("Helvetica", 10)
            c.drawString(50 * mm, height - 20 * mm, f"{title}（第 {page_num} 页 - 空白）")
            c.showPage()
            continue
        else:
            c.setFont("Helvetica", 10)
            c.drawString(50 * mm, height - 20 * mm, f"{title}（第 {page_num} 页）")
            y = height - 35 * mm

        c.setFont("Helvetica", 12)
        # 渲染当前页的段落，按段落数限，不按行数限
        para_on_page = 0
        while para_index < total_paragraphs and para_on_page < paras_per_page:
            text = content[para_index]

            # 自动换行渲染
            if c.stringWidth(text, "Helvetica", 12) > width - 60 * mm:
                words = text.split()
                line = ""
                for word in words:
                    test_line = f"{line} {word}".strip()
                    if c.stringWidth(test_line, "Helvetica", 12) > width - 60 * mm:
                        c.drawString(30 * mm, y, line)
                        y -= LINE_HEIGHT
                        line = word
                    else:
                        line = test_line
                if line:
                    c.drawString(30 * mm, y, line)
                    y -= LINE_HEIGHT
            else:
                c.drawString(30 * mm, y, text)
                y -= LINE_HEIGHT

            para_index += 1
            para_on_page += 1

        # 表格（最后一页）
        if para_index >= total_paragraphs:
            table = data.get("table")
            if table and y > 40 * mm:
                y -= 10 * mm
                headers = table.get("headers", [])
                rows = table.get("rows", [])
                c.setFont("Helvetica-Bold", 10)
                x = 30 * mm
                for h in headers:
                    c.drawString(x, y, h)
                    x += 40 * mm
                y -= 8 * mm
                c.setFont("Helvetica", 10)
                for row in rows:
                    if y < 20 * mm:
                        break
                    x = 30 * mm
                    for cell in row:
                        c.drawString(x, y, str(cell))
                        x += 40 * mm
                    y -= 8 * mm

        if para_index < total_paragraphs or (page_count > 0 and page_num < page_count):
            c.showPage()

    c.save()
    extra = "，已加密" if encrypt else ""
    print(f"[OK] PDF 文件已生成: {output_path}，共 {page_num} 页{extra}")


# ============================================================
# 主入口
# ============================================================

def _adjust_content_for_size(data: dict, fmt: str) -> dict:
    """根据 file_size 调整内容量，使生成的文件接近目标大小

    校准值（基于实际测试）：
      - Word(docx): ~200 bytes/paragraph (40 词随机文本), base ~10KB
      - Excel(xlsx): ~200 bytes/row (8 列), base ~8KB
      - PPT(pptx): ~1000 bytes/slide, base ~15KB
      - MD: target 字节直接写入纯文本
      - PDF: ~200 bytes/paragraph, base ~5KB
    """
    target = data.get("file_size", 0)
    if not target or target <= 0:
        return data

    data = dict(data)  # 不修改原数据

    # ── Markdown（纯文本，直接计算） ──
    if fmt == "md":
        # 每个 section 约 910 字节（65 字符 text * 5 = 325 字符，UTF-8 编码）
        BYTES_PER_SECTION = 910
        sections_needed = max(1, min(10000, int(target / BYTES_PER_SECTION)))
        text = ("这是用于测试文件大小的自动生成内容。请忽略此内容的实际含义。"
                "本文件由 Python 脚本自动生成，用于测试文件下载和文档生成功能。")
        sections = []
        for i in range(sections_needed):
            sections.append({
                "heading": f"章节 {i + 1}",
                "content": text * 5,
            })
        data["sections"] = sections
        data.setdefault("title", "示例文档")
        data.setdefault("date", datetime.now().strftime("%Y-%m-%d"))
        data.setdefault("author", "系统")
        return data

    # ── 压缩格式 ──
    BPU = {"word": 200, "excel": 60, "ppt": 1000, "pdf": 200}
    OVERHEAD = {"word": 10000, "excel": 8000, "ppt": 15000, "pdf": 5000}
    # 各格式最大单位数（防止生成太慢）
    MAX_UNITS = {"word": 50000, "excel": 400000, "ppt": 500, "pdf": 50000}

    bpu = BPU.get(fmt, 200)
    overhead = OVERHEAD.get(fmt, 5000)
    max_units = MAX_UNITS.get(fmt, 5000)

    # 计算需要的单位数
    units = max(1, min(max_units, int((target - overhead) / bpu)))

    # ── Word ──
    if fmt == "word":
        random.seed(42)
        paragraphs = []
        for _ in range(units):
            text = ' '.join(''.join(random.choices(string.ascii_lowercase, k=random.randint(3, 10)))
                          for _ in range(40))
            paragraphs.append({"text": text, "bold": False})
        data["paragraphs"] = paragraphs
        data.setdefault("title", "示例文档")

    # ── Excel ──
    elif fmt == "excel":
        # 用户显式设置了 row_count 或 column_count 时，不再覆盖
        if "row_count" not in data or data.get("row_count", 0) <= 0:
            data["row_count"] = max(10, units)
        if "column_count" not in data:
            data.setdefault("columns", ["姓名", "年龄", "部门", "薪资", "手机号", "邮箱", "入职日期", "状态"])

    # ── PPT ──
    elif fmt == "ppt":
        slides = []
        for i in range(min(units, 5000)):
            slides.append({
                "title": f"第 {i + 1} 页",
                "items": [f"内容项 {j + 1}" for j in range(5)],
            })
        data["slides"] = slides
        data.setdefault("title", "演示文稿")
        data.setdefault("subtitle", "自动生成")

    # ── PDF ──
    elif fmt == "pdf":
        data.setdefault("title", "示例文档")
        if _as_bool(data.get("empty_content")):
            data["content"] = []
            return data
        random.seed(42)
        page_count = 0 if _as_bool(data.get("empty_page_count")) else data.get("page_count", 0)
        if page_count > 0:
            units = max(units, page_count * 15)
        paragraphs = []
        for _ in range(units):
            text = ' '.join(''.join(random.choices(string.ascii_lowercase, k=random.randint(3, 10)))
                          for _ in range(40))
            paragraphs.append(text)
        data["content"] = paragraphs

    # ── 图片 ──
    elif fmt == "image":
        random.seed(42)
        # 图片通过散点图（大量数据点）+ figsize + dpi 控制文件大小
        # 校准: 散点图 PNG 大小 ≈ 数据点 × 0.05KB + figsize_factor × 0.1KB
        # 100KB → 3000 pts, 10×6 @ 150dpi
        # 500KB → 15000 pts, 14×8 @ 200dpi
        # 1MB → 30000 pts, 16×10 @ 250dpi
        # 2MB → 60000 pts, 20×12 @ 300dpi
        # 5MB → 100000 pts, 24×16 @ 300dpi
        # 10MB → 200000 pts, 28×18 @ 350dpi
        # 15MB → 200000 pts, 30×20 @ 400dpi
        # 20MB → 200000 pts, 32×22 @ 450dpi
        FIGURE_MAP = [
            (102400,     3000,  (10, 6),   150),
            (512000,     15000, (14, 8),   200),
            (1048576,    30000, (16, 10),  250),
            (2097152,    40000, (20, 12),  300),
            (5242880,    100000, (24, 16), 300),
            (10485760,   150000, (26, 16), 300),
            (15728640,   200000, (28, 18), 300),
            (20971520,   200000, (30, 20), 350),
        ]
        config = FIGURE_MAP[0]
        for limit, pts, size, d in FIGURE_MAP:
            if target <= limit:
                config = (pts, size, d)
                break
            config = (pts, size, d)
        n_points, figsize, dpi = config
        data["chart_type"] = "scatter"
        data["figsize"] = figsize
        data["dpi"] = dpi
        data["n_points"] = n_points
        data.setdefault("title", "测试图表")

    return data


GENERATORS = {
    "word": lambda data, path: generate_word(
        data if data.get("title") or data.get("paragraphs") else {
            "title": "示例文档",
            "paragraphs": ["这是自动生成的 Word 文档。", "段落内容由 Python 脚本自动填充。"],
            "table": {"headers": ["项目", "描述", "状态"], "rows": [["模块A", "自动生成", "已完成"], ["模块B", "待处理", "进行中"]]}
        }, path
    ),
    "excel": lambda data, path: generate_excel_mock({
        "columns": data.get("columns") or ["姓名", "年龄", "部门", "薪资", "手机号", "邮箱", "入职日期", "状态"],
        "column_count": data.get("column_count", 0),
        "row_count": data.get("row_count", 20),
        "row_size": data.get("row_size", 10),
        "file_count": data.get("file_count", 1),
        "output_dir": data.get("output_dir", ""),
        "sheet_name": data.get("sheet_name", "Sheet1"),
    }, path),
    "ppt": lambda data, path: generate_ppt(
        data if data.get("slides") else {
            "title": "演示文稿",
            "subtitle": "自动生成",
            "slides": [
                {"title": "第一页", "items": ["内容点一", "内容点二", "内容点三"]},
                {"title": "第二页", "items": ["核心功能", "技术架构", "部署方案"]},
            ]
        }, path
    ),
    "image": lambda data, path: generate_image(
        data if data.get("chart_type") else {
            "chart_type": "bar",
            "title": "柱状图示例",
            "categories": ["1月", "2月", "3月", "4月"],
            "values": [100, 150, 130, 180],
        }, path
    ),
    "md": lambda data, path: generate_markdown(
        data if data.get("sections") else {
            "title": "示例文档",
            "date": datetime.now().strftime("%Y-%m-%d"),
            "author": "系统",
            "sections": [
                {"heading": "第一章", "content": "这是第一章的内容，由 Python 脚本自动生成。"},
                {"heading": "第二章", "content": "这是第二章的内容。"},
            ]
        }, path
    ),
    "pdf": lambda data, path: generate_pdf(data, path),
}


def main():
    if len(sys.argv) < 4:
        print("用法: python3 generate_doc.py <format> <data.json> <output_path>")
        print("支持的格式: " + ", ".join(GENERATORS.keys()))
        sys.exit(1)

    fmt = sys.argv[1].lower()
    data_path = sys.argv[2]
    output_path = sys.argv[3]

    if fmt not in GENERATORS:
        print(f"不支持的格式: {fmt}")
        print(f"支持的格式: {', '.join(GENERATORS.keys())}")
        sys.exit(1)

    if not os.path.exists(data_path):
        print(f"数据文件不存在: {data_path}")
        sys.exit(1)

    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # 根据 file_size 调整内容量
    data = _adjust_content_for_size(data, fmt)

    # 如果 data 中指定了 output_dir，则保存到该目录
    output_dir = data.get("output_dir", "")
    if output_dir:
        output_dir = os.path.abspath(output_dir)
        os.makedirs(output_dir, exist_ok=True)
        # 将 output_path 解析到 output_dir 下
        output_path = os.path.join(output_dir, os.path.basename(output_path))

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    file_count = data.get("file_count", 1)
    output_dir = data.get("output_dir", "")

    if file_count > 1 and fmt != "excel":
        # 非 Excel 格式：生成多份拷贝
        base = os.path.splitext(output_path)[0]
        base_name = os.path.basename(base)
        ext = os.path.splitext(output_path)[1]
        total_size = 0
        for i in range(file_count):
            part_name = f"{base_name}_{i + 1}{ext}"
            if output_dir:
                part_path = os.path.join(output_dir, part_name)
            else:
                part_path = f"{base}_{i + 1}{ext}"
            # 第一个文件也写入原始 output_path，浏览器下载时能读到
            if i == 0 and not output_dir:
                GENERATORS[fmt](data, output_path)
            GENERATORS[fmt](data, part_path)
            size = os.path.getsize(part_path)
            total_size += size
            print(f"  [{i + 1}/{file_count}] {os.path.basename(part_path)} ({size} bytes)")
        print(f"[OK] {fmt.upper()} {file_count} 个文件已生成，总大小: {total_size} bytes")
    else:
        # Excel（内部处理 file_count）或单文件
        GENERATORS[fmt](data, output_path)


if __name__ == "__main__":
    main()